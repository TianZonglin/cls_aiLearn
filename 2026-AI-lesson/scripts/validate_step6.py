import csv
import json
import shutil
import tempfile
from pathlib import Path
from uuid import uuid4

from docx import Document as DocxDocument
from fastapi.testclient import TestClient
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation
from reportlab.pdfgen import canvas
import xlwt

from app.core.db import SessionLocal, init_db
from app.core.paths import ensure_runtime_dirs
from app.main import app
from app.models.document_chunks import DocumentChunk
from app.models.documents import Document
from app.repositories.knowledge_bases import utc_now
from app.services.document_parsers import CC_CONVERTER


def print_result(name: str, passed: bool, detail: str) -> None:
    status = "PASS" if passed else "FAIL"
    print(f"[{status}] {name}: {detail}")


def create_pdf(path: Path) -> None:
    packet = canvas.Canvas(str(path))
    packet.drawString(72, 720, "PDF 中文样例 内容测试")
    packet.save()


def create_docx(path: Path) -> None:
    doc = DocxDocument()
    doc.add_paragraph("DOCX 中文样例 内容测试")
    doc.save(str(path))


def create_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    textbox = slide.shapes.add_textbox(left=0, top=0, width=6000000, height=1000000)
    textbox.text = "PPTX 中文样例 内容测试"
    prs.save(str(path))


def create_xlsx(path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "课程表"
    ws["A1"] = "字段"
    ws["B1"] = "值"
    ws["A2"] = "主题"
    ws["B2"] = "XLSX 中文样例 内容测试"
    wb.save(str(path))


def create_xls(path: Path) -> None:
    wb = xlwt.Workbook()
    ws = wb.add_sheet("旧表")
    ws.write(0, 0, "字段")
    ws.write(0, 1, "值")
    ws.write(1, 0, "主题")
    ws.write(1, 1, "XLS 中文样例 内容测试")
    wb.save(str(path))


def create_csv(path: Path) -> None:
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerow(["字段", "值"])
        writer.writerow(["主题", "CSV 中文样例 内容测试"])


def create_png(path: Path) -> None:
    image = Image.new("RGB", (640, 240), color="white")
    draw = ImageDraw.Draw(image)
    draw.text((20, 100), "網絡 SAMPLE 123", fill="black")
    image.save(str(path))


def main() -> int:
    ensure_runtime_dirs()
    init_db()

    files_root = Path(__file__).resolve().parents[1] / "storage" / "files"
    if files_root.exists():
        shutil.rmtree(files_root)

    temp_dir = Path(tempfile.mkdtemp(prefix="step6_"))
    try:
        pdf_path = temp_dir / "sample.pdf"
        docx_path = temp_dir / "sample.docx"
        pptx_path = temp_dir / "sample.pptx"
        xlsx_path = temp_dir / "sample.xlsx"
        xls_path = temp_dir / "sample.xls"
        csv_path = temp_dir / "sample.csv"
        png_path = temp_dir / "sample.png"

        create_pdf(pdf_path)
        create_docx(docx_path)
        create_pptx(pptx_path)
        create_xlsx(xlsx_path)
        create_xls(xls_path)
        create_csv(csv_path)
        create_png(png_path)

        client = TestClient(app)
        kb_response = client.post(
            "/knowledge-bases",
            json={"name": "Step6 Parse KB", "description": "for parsing validation"},
        )
        kb_response.raise_for_status()
        knowledge_base_id = kb_response.json()["id"]

        upload_files = [
            ("files", ("sample.pdf", pdf_path.read_bytes(), "application/pdf")),
            ("files", ("sample.docx", docx_path.read_bytes(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")),
            ("files", ("sample.pptx", pptx_path.read_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
            ("files", ("sample.xlsx", xlsx_path.read_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")),
            ("files", ("sample.xls", xls_path.read_bytes(), "application/vnd.ms-excel")),
            ("files", ("sample.csv", csv_path.read_bytes(), "text/csv")),
            ("files", ("sample.png", png_path.read_bytes(), "image/png")),
        ]
        upload_response = client.post(
            "/documents/upload",
            data={"knowledge_base_id": knowledge_base_id},
            files=upload_files,
        )
        upload_response.raise_for_status()
        upload_data = upload_response.json()

        now = utc_now()
        html_document_id = str(uuid4())
        html_document = Document(
            id=html_document_id,
            knowledge_base_id=knowledge_base_id,
            name="Step6 HTML Page",
            source_type="url",
            file_type="html",
            mime_type="text/html",
            source_url="https://example.com/step6",
            original_path=None,
            storage_path="https://example.com/step6",
            file_size=64,
            checksum=None,
            parse_status="pending",
            parse_error=None,
            preview_text="这是一个 HTML 中文内容样例，用于解析与切片验证。",
            summary_text=None,
            page_count=None,
            retry_count=0,
            last_parsed_at=None,
            metadata_json=None,
            created_at=now,
            updated_at=now,
        )

        session = SessionLocal()
        try:
            session.add(html_document)
            session.commit()
        finally:
            session.close()

        all_document_ids = [item["document_id"] for item in upload_data["success"]] + [html_document_id]

        parse_results = []
        for document_id in all_document_ids:
            result = client.post(f"/documents/{document_id}/index")
            parse_results.append((document_id, result.status_code, result.json()))

        done_count = sum(1 for _, status_code, payload in parse_results if status_code == 200 and payload["parse_status"] == "done")
        print_result("多类型解析到 done", done_count == len(all_document_ids), json.dumps(parse_results, ensure_ascii=False))

        docs_response = client.get(f"/documents?knowledge_base_id={knowledge_base_id}")
        docs_response.raise_for_status()
        docs = docs_response.json()
        preview_ok = all((doc["preview_text"] or "").strip() for doc in docs if doc["parse_status"] == "done")
        print_result("preview_text 非空", preview_ok, json.dumps([(doc["name"], doc["parse_status"], doc["preview_text"][:40] if doc["preview_text"] else "") for doc in docs], ensure_ascii=False))

        session = SessionLocal()
        try:
            chunk_count = session.query(DocumentChunk).count()
            print_result("document_chunks 已生成", chunk_count > 0, f"chunk_count={chunk_count}")
        finally:
            session.close()

        retry_upload = client.post(
            "/documents/upload",
            data={"knowledge_base_id": knowledge_base_id},
            files=[("files", ("retry.csv", "字段,值\n主题,重试解析中文内容".encode("utf-8"), "text/csv"))],
        )
        retry_upload.raise_for_status()
        failed_document_id = retry_upload.json()["success"][0]["document_id"]

        session = SessionLocal()
        try:
            failed_document = session.get(Document, failed_document_id)
            if failed_document is not None:
                failed_document.parse_status = "failed"
                failed_document.parse_error = "manual failure for retry"
                session.add(failed_document)
                session.commit()
        finally:
            session.close()

        retry_result = client.post(f"/documents/{failed_document_id}/retry-parse")
        retry_ok = retry_result is not None and retry_result.status_code == 200 and retry_result.json()["parse_status"] == "done"
        print_result("失败后可重试解析", retry_ok, retry_result.text if retry_result is not None else "retry document missing")

        chinese_ok = any("中文" in (doc["preview_text"] or "") for doc in docs if doc["parse_status"] == "done")
        print_result("中文预览不乱码", chinese_ok, json.dumps([(doc["name"], doc["preview_text"][:60] if doc["preview_text"] else "") for doc in docs], ensure_ascii=False))

        xls_ok = any(doc["name"] == "sample.xls" and doc["parse_status"] == "done" for doc in docs)
        print_result("XLS 旧格式解析", xls_ok, json.dumps([(doc["name"], doc["parse_status"]) for doc in docs], ensure_ascii=False))

        ocr_doc = next((doc for doc in docs if doc["name"] == "sample.png"), None)
        ocr_preview_ok = ocr_doc is not None and bool((ocr_doc["preview_text"] or "").strip())
        print_result("OCR 图片解析链路", ocr_preview_ok, json.dumps(ocr_doc, ensure_ascii=False) if ocr_doc else "missing OCR doc")

        normalization_ok = CC_CONVERTER.convert("網絡") == "网络"
        print_result("OCR 繁转简规范化", normalization_ok, CC_CONVERTER.convert("網絡"))
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
