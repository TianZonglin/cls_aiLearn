import csv
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional
from uuid import uuid4

from docx import Document as DocxDocument
from opencc import OpenCC
from openpyxl import load_workbook
from PIL import Image, ImageOps
from pptx import Presentation
from pypdf import PdfReader
from rapidocr_onnxruntime import RapidOCR
import xlrd

from app.models.document_chunks import DocumentChunk
from app.models.documents import Document
from app.repositories.knowledge_bases import utc_now


OCR_ENGINE: Optional[RapidOCR] = None
CC_CONVERTER = OpenCC("t2s")


@dataclass
class ParsedSegment:
    text: str
    location_label: str
    page_number: Optional[int] = None
    sheet_name: Optional[str] = None


def get_ocr_engine() -> RapidOCR:
    global OCR_ENGINE
    if OCR_ENGINE is None:
        OCR_ENGINE = RapidOCR()
    return OCR_ENGINE


def load_text_from_document(document: Document, project_root: Path) -> List[ParsedSegment]:
    if document.source_type == "url":
        return [ParsedSegment(text=document.preview_text or "", location_label="URL Content")]

    file_path = project_root / document.storage_path
    file_type = document.file_type

    if file_type == "pdf":
        return parse_pdf(file_path)
    if file_type == "docx":
        return parse_docx(file_path)
    if file_type == "pptx":
        return parse_pptx(file_path)
    if file_type == "excel":
        return parse_excel(file_path)
    if file_type == "csv":
        return parse_csv(file_path)
    if file_type in {"png", "jpg", "jpeg"}:
        return parse_image_ocr(file_path)

    raise ValueError(f"Unsupported parser for file_type={file_type}")


def parse_pdf(file_path: Path) -> List[ParsedSegment]:
    reader = PdfReader(str(file_path))
    segments: List[ParsedSegment] = []
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            segments.append(ParsedSegment(text=text, location_label=f"Page {index}", page_number=index))
    if not segments:
        raise ValueError("PDF text extraction returned empty content.")
    return segments


def parse_docx(file_path: Path) -> List[ParsedSegment]:
    doc = DocxDocument(str(file_path))
    paragraphs = [paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip()]
    text = "\n".join(paragraphs).strip()
    if not text:
        raise ValueError("DOCX text extraction returned empty content.")
    return [ParsedSegment(text=text, location_label="Document Body")]


def parse_pptx(file_path: Path) -> List[ParsedSegment]:
    prs = Presentation(str(file_path))
    segments: List[ParsedSegment] = []
    for index, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                content = shape.text.strip()
                if content:
                    parts.append(content)
        text = "\n".join(parts).strip()
        if text:
            segments.append(ParsedSegment(text=text, location_label=f"Slide {index}", page_number=index))
    if not segments:
        raise ValueError("PPTX text extraction returned empty content.")
    return segments


def parse_excel(file_path: Path) -> List[ParsedSegment]:
    if file_path.suffix.lower() == ".xls":
        return parse_xls(file_path)

    workbook = load_workbook(filename=str(file_path), read_only=True, data_only=True)
    segments: List[ParsedSegment] = []
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if values:
                rows.append(" | ".join(values))
        text = "\n".join(rows).strip()
        if text:
            segments.append(ParsedSegment(text=text, location_label=f"Sheet {sheet_name}", sheet_name=sheet_name))
    if not segments:
        raise ValueError("Excel text extraction returned empty content.")
    return segments


def parse_xls(file_path: Path) -> List[ParsedSegment]:
    workbook = xlrd.open_workbook(str(file_path))
    segments: List[ParsedSegment] = []
    for sheet in workbook.sheets():
        rows = []
        for row_index in range(sheet.nrows):
            values = [str(cell).strip() for cell in sheet.row_values(row_index) if str(cell).strip()]
            if values:
                rows.append(" | ".join(values))
        text = "\n".join(rows).strip()
        if text:
            segments.append(ParsedSegment(text=text, location_label=f"Sheet {sheet.name}", sheet_name=sheet.name))
    if not segments:
        raise ValueError("XLS text extraction returned empty content.")
    return segments


def parse_csv(file_path: Path) -> List[ParsedSegment]:
    rows = []
    with file_path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
        reader = csv.reader(handle)
        for row in reader:
            values = [cell.strip() for cell in row if cell and cell.strip()]
            if values:
                rows.append(" | ".join(values))
    text = "\n".join(rows).strip()
    if not text:
        raise ValueError("CSV text extraction returned empty content.")
    return [ParsedSegment(text=text, location_label="CSV Rows")]


def parse_image_ocr(file_path: Path) -> List[ParsedSegment]:
    image = Image.open(file_path)
    parts = []

    for candidate in build_ocr_candidates(image):
        result, _ = get_ocr_engine()(candidate)
        if result:
            current_parts = []
            for item in result:
                if len(item) >= 2 and item[1]:
                    current_parts.append(CC_CONVERTER.convert(str(item[1])))
            if current_parts:
                parts = current_parts
                break

    text = "\n".join(parts).strip()
    if not text:
        raise ValueError("OCR returned empty content.")
    return [ParsedSegment(text=text, location_label="OCR Text")]


def build_ocr_candidates(image: Image.Image) -> List[Image.Image]:
    base = ImageOps.exif_transpose(image).convert("RGB")
    gray = ImageOps.grayscale(base)
    enlarged = gray.resize((gray.width * 2, gray.height * 2))
    high_contrast = ImageOps.autocontrast(enlarged)
    binary = high_contrast.point(lambda value: 0 if value < 180 else 255, mode="1").convert("RGB")

    return [
        base,
        enlarged.convert("RGB"),
        high_contrast.convert("RGB"),
        binary,
    ]


def chunk_segments(document: Document, segments: List[ParsedSegment]) -> List[DocumentChunk]:
    now = utc_now()
    chunks: List[DocumentChunk] = []
    chunk_index = 0

    for segment in segments:
        text = segment.text.strip()
        if not text:
            continue
        start = 0
        window = 800
        overlap = 120
        while start < len(text):
            end = min(len(text), start + window)
            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append(
                    DocumentChunk(
                        id=str(uuid4()),
                        document_id=document.id,
                        knowledge_base_id=document.knowledge_base_id,
                        chunk_index=chunk_index,
                        text=chunk_text,
                        token_count=len(chunk_text),
                        location_label=segment.location_label,
                        page_number=segment.page_number,
                        sheet_name=segment.sheet_name,
                        start_offset=start,
                        end_offset=end,
                        vector_id=f"chunk-{document.id}-{chunk_index}",
                        created_at=now,
                    )
                )
                chunk_index += 1
            if end >= len(text):
                break
            start = max(end - overlap, start + 1)

    if not chunks:
        raise ValueError("Chunking produced no chunks.")
    return chunks


def build_preview_text(segments: List[ParsedSegment]) -> str:
    full_text = "\n".join(segment.text for segment in segments if segment.text).strip()
    return full_text[:1000]


def build_summary_placeholder(segments: List[ParsedSegment]) -> str:
    full_text = "\n".join(segment.text for segment in segments if segment.text).strip()
    if not full_text:
        return ""
    return full_text[:180]
